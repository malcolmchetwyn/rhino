
import os
from typing import List, Optional
from fastapi import FastAPI, Request, Form, Depends, UploadFile
from fastapi.responses import HTMLResponse, StreamingResponse, FileResponse
from fastapi import HTTPException
from fastapi.templating import Jinja2Templates
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
import openai
from langchain.vectorstores import FAISS
from openai.error import OpenAIError
from dotenv import load_dotenv
import tiktoken
from concurrent.futures import ThreadPoolExecutor
from langchain.embeddings import HuggingFaceBgeEmbeddings
import re
import logging
from send_email import send_email 
from fastapi.responses import JSONResponse 
from vector_news import get_news_details, preload_news_indices
from vector_all_sf_user_guides_data import preload_sf_user_guides_indices, get_sf_user_guides_details
from fastapi import FastAPI, Form, Depends, HTTPException, Request, Response, status
from passlib.context import CryptContext
from starlette.templating import Jinja2Templates
from starlette.staticfiles import StaticFiles
import redis
import os
import random
import string
import smtplib
from email.message import EmailMessage
from fastapi.responses import RedirectResponse
from dotenv import load_dotenv
from send_email import send_verification_email, send_email
from fastapi import Cookie, HTTPException
from datetime import timedelta
from docx import Document
import uuid
from PyPDF2 import PdfReader
from openpyxl import load_workbook
from pptx import Presentation
import sys
from fastapi import UploadFile, File, Form
import shutil
from pathlib import Path
import os
import shutil
import uuid
from werkzeug.utils import secure_filename  
from pathlib import Path 
from fastapi import FastAPI, UploadFile, Form

UPLOAD_DIR = "uploads"
Path(UPLOAD_DIR).mkdir(parents=True, exist_ok=True)
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'xlsx', 'docx', 'pptx', 'doc'}
os.makedirs("uploaded_files", exist_ok=True)

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
r = redis.Redis(host='localhost', port=6379, db=0)
CREDENTIALS_FILE = "credentials.txt"

load_dotenv()
loaded_indexes = {}

#loading vector stores.
preload_news_indices()
#preload_sf_user_guides_indices()

load_dotenv()  # Load the environment variables from .env file

# Load your HuggingFaceBgeEmbeddings
embeddings = HuggingFaceBgeEmbeddings(model_name="BAAI/bge-large-en-v1.5")
 
app = FastAPI()
app.mount("/static", StaticFiles(directory="static"), name="static")

templates = Jinja2Templates(directory="templates")

openai.api_key = os.getenv('OPENAI_API_KEY')

# Define a custom log format
log_format = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')

# Create a handler to log to the console
console_handler = logging.StreamHandler()
console_handler.setFormatter(log_format)

# Create a handler to log to the file
file_handler = logging.FileHandler('post_requests.log')
file_handler.setFormatter(log_format)

# Get the root logger
root_logger = logging.getLogger()
root_logger.setLevel(logging.INFO) # or the lowest log level you want to capture
root_logger.addHandler(console_handler)
root_logger.addHandler(file_handler)


def send_text( 
    prompt=None,
    text_data=None,
    chat_model="gpt-4-1106-preview", 
    model_token_limit=120000,
    max_tokens=120000,
):
    """
    Send the prompt at the start of the conversation and then send chunks of text_data to ChatGPT via the OpenAI API.
    If the text_data is too long, it splits it into chunks and sends each chunk separately.
    """

    # Check if the necessary arguments are provided
    if not prompt:
        print("Error: Missing prompt.")
        return "Error: Prompt is missing. Please provide a prompt."
    if not text_data:
        print("Error: Missing text data.")
        return "Error: Text data is missing. Please provide some text data."

    # Initialize the tokenizer
    tokenizer = tiktoken.encoding_for_model("gpt-4-1106-preview")
    print("Tokenizer initialized.")

    # Encode the text_data into token integers
    token_integers = tokenizer.encode(text_data)
    print(f"Text data encoded into {len(token_integers)} tokens.")

    # Split the token integers into chunks based on max_tokens
    chunk_size = max_tokens - len(tokenizer.encode(prompt))
    chunks = [
        token_integers[i : i + chunk_size]
        for i in range(0, len(token_integers), chunk_size)
    ]
    print(f"Text data split into {len(chunks)} chunks.")

    # Decode token chunks back to strings
    chunks = [tokenizer.decode(chunk) for chunk in chunks]

    responses = []
    messages = [
        {"role": "user", "content": prompt},
        {
            "role": "user",
            "content": "To provide the context for the above prompt, I will send you text in parts. When I am finished, I will tell you 'ALL PARTS SENT'. Do not answer until you have received all the parts.",
        },
    ]

    for index, chunk in enumerate(chunks):
        print(f"Processing chunk {index + 1}/{len(chunks)}...")
        
        messages.append({"role": "user", "content": chunk})

        # Check if total tokens exceed the model's limit and remove oldest chunks if necessary
        while (
            sum(len(tokenizer.encode(msg["content"])) for msg in messages)
            > model_token_limit
        ):
            messages.pop(1)  # Remove the oldest chunk
            print("Removed oldest chunk due to token limit.")

        response = openai.ChatCompletion.create(model=chat_model, messages=messages)
        chatgpt_response = response.choices[0].message["content"].strip()
        responses.append(chatgpt_response)
        print(f"Received response for chunk {index + 1}.")

    # Add the final "ALL PARTS SENT" message
    print("Sending 'ALL PARTS SENT' message...")
    messages.append({"role": "user", "content": "ALL PARTS SENT"})
    response = openai.ChatCompletion.create(model=chat_model, messages=messages)
    final_response = response.choices[0].message["content"].strip()
    responses.append(final_response)

    print("Processing completed.")
    return responses


@app.get("/")
def read_root(request: Request, session_token: str = Cookie(None)):
    # If the user is already authenticated, redirect to the main page
    if session_token and r.exists(session_token):
        return RedirectResponse(url="/chat", status_code=status.HTTP_303_SEE_OTHER)
    # Otherwise, show the login/signup page
    return templates.TemplateResponse("index.html", {"request": request})
  

def num_tokens_from_string(string: str, encoding_name: str) -> int:
    """Returns the number of tokens in a text string."""
    encoding = tiktoken.get_encoding(encoding_name)
    num_tokens = len(encoding.encode(string))
    return num_tokens

class Message(BaseModel): 
    role: str
    content: str

async def generate(messages: List[Message], model_type: str):
    try:
        for message in messages:
            print(message)
        print("trying to connect")
        
        response = await openai.ChatCompletion.acreate(
            model=model_type,
            messages=[message.dict() for message in messages],
            stream=True
        )
       
        async for chunk in response:
            content = chunk['choices'][0]['delta'].get('content', '')
            if content:
                yield content

    except OpenAIError as e:
        print(f"{type(e).__name__}: {str(e)}")
        yield f"{type(e).__name__}: {str(e)}"


def extract_text_from_excel(file_path):
    workbook = load_workbook(filename=file_path)
    text = ""
    for sheet in workbook:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    text += str(cell) + "\t"
            text += "\n"
    return text

def extract_text_from_powerpoint(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_file(file_path, file_type):
    if file_type == '.pdf':
        with open(file_path, "rb") as pdf_file:
            pdf = PdfReader(pdf_file)
            text = ""
            for page in pdf.pages:
                text += page.extract_text()

            if not text.strip():
                raise ValueError("Failed to extract text from PDF")
            return text

    elif file_type == '.txt':
        with open(file_path, "r") as file:
            content = file.read()
            return content
 
    elif file_type == '.docx':
        doc = Document(file_path)
        content = "\n".join([p.text for p in doc.paragraphs if p.text.strip() != ""])
        return content
    
    elif file_type == '.xlsx':
        return extract_text_from_excel(file_path)

    elif file_type == '.pptx':
        return extract_text_from_powerpoint(file_path)


    else:
        raise ValueError(f"Unsupported file type: {file_type}")


# Note: I've removed model_type from the request model
class Gpt4Request(BaseModel):
    messages: List[Message]

from time import sleep

def error_stream(error_message: str):
    for word in error_message.split():
        yield word + " "  # adding a space after each word
        sleep(0.05)  # you can adjust this delay as needed




def custom_response_stream(error_message: str):
    words = error_message.split(' ')
    for word in words:
        if '\n' in word:
            # split at newline and yield words and newline separately
            parts = word.split('\n')
            for i, part in enumerate(parts):
                if i != 0:  # not the first part, so we yield a newline first
                    yield '\n'
                yield part + " "
        else:
            yield word + " "
        sleep(0.5)  # you can adjust this delay as needed

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def clean_pdf_text(text):
    # Split the data into lines and remove redundant spaces
    lines = [line.strip() for line in text.split("\n") if line.strip() != ""]

    enhanced_text = []

    for idx, line in enumerate(lines):
        # Financial year data followed by metrics, e.g. "5% on FY22$10,164m"
        if re.match(r"^\d+% on FY\d+\$\d+m", line):
            split_line = re.split(r"(\$)", line, 1)
            enhanced_text.append(split_line[0] + split_line[1] + split_line[2])
            enhanced_text.append("\n")  # Add an extra line for clarity

        # Financial metrics followed by details  
        elif re.match(r"\$\d+\.?\d*m", line) or re.match(r"\d+\.?\d*%", line):
            enhanced_text.append(line)
            if (idx + 1 < len(lines)):
                enhanced_text.append(lines[idx + 1])
            enhanced_text.append("\n")  # Add an extra line for clarity

        # Titles and sub-titles
        elif line in ["Balance Sheet", "Statement of Comprehensive Income", "RACV Annual Report",
                      "Investments Accounted for Using the Equity Method", "Balance Sheet", "COMMONWEAL TH BANK"]:
            enhanced_text.append("\n\n" + line + "\n" + "-"*len(line) + "\n")

        # Bullet points (e.g., • Arevo Pty Ltd)
        elif re.match(r"•\s", line):
            enhanced_text.append("  - " + line.split('•')[1])

        # Special notes
        elif line.startswith("1  Refer to note"):
            enhanced_text.append("\nNote: " + line + "\n")

        # Categories followed by descriptions (e.g., "Strategic", "Financial")
        elif line in ["Strategic", "Financial", "Non ‑financial", "EmergingRisk", "Financial risk"]:
            enhanced_text.append(line + ":")
            if (idx + 1 < len(lines)):
                enhanced_text.append("  " + lines[idx + 1])
            enhanced_text.append("\n")

        # Page numbers
        elif re.match(r"^\d{1,3} \d{1,3}$", line):
            continue

        else:
            enhanced_text.append(line)

    # Concatenate all the enhanced text into a single string
    cleaned_file = "\n".join(enhanced_text)

    return cleaned_file



import requests
import io
import json


@app.post("/gpt4")
async def gpt4(messages: str = Form(...), file: UploadFile = None): 
    #raise Exception("429 Rate limit reached for requests")


    message_data = json.loads(messages)

    user_request = message_data[-1]["content"]

    logging.info(f"PROMPT: {user_request}")

    if not user_request.strip().split():
            error_message= "Hold your horses, keyboard cowboy! 🐎💨 It seems you've sent an invisible ink message. I'll need at least one word to wrangle up some answers. Try again with a bit of text, partner!"
            logging.error(f"NO PROMPT: User Entered No Prompt")
            return StreamingResponse(error_stream(error_message), media_type='text/plain')
    
        # Check if the user request contains the placeholder text
    if "[INSERT COMPANY NAME]" in user_request:
        error_message= "Whoops! Looks like you forgot to insert your company name 😄. Please replace '[INSERT COMPANY NAME]' with your real company name and try again."
        logging.error(f"NO COMPANY: User Didnt enter a compoany name")
        return StreamingResponse(error_stream(error_message), media_type='text/plain')

    
    try:
        cleaned_file = ""
        # Save the uploaded file to the 'uploads' directory
      
        if file:

            if allowed_file(file.filename):

                filename = secure_filename(file.filename)
                
                # Ensure the directory exists; if not, create it
                if not os.path.exists(UPLOAD_DIR):
                    print(f"Creating directory: {app.config['UPLOAD_FOLDER']}")
                    try:
                        os.makedirs(UPLOAD_DIR)
                        print(f"Directory {UPLOAD_DIR} created successfully!")
                    except Exception as dir_error:
                        print(f"Error creating directory: {dir_error}")
                        return f"Error creating directory: {dir_error}"

                # Attempt to save the file
                try:
                    # Generate a random UUID and take the first 8 characters for brevity
                    random_string = str(uuid.uuid4())[:8]

                    # Add the random string to the filename
                    base, ext = os.path.splitext(filename)
                    randomized_filename = f"{base}_{random_string}{ext}"

                    file_path = os.path.join(UPLOAD_DIR, randomized_filename)
                    print(f"Attempting to save file to: {file_path}")  # Print the file path for debugging
                    logging.info(f"Attempting to save file to: {file_path}")

                    if os.access(os.path.dirname(file_path), os.W_OK):
                        with open(file_path, "wb") as buffer:
                            shutil.copyfileobj(file.file, buffer)
                        print(f"File saved successfully to {file_path}!")
                        logging.info(f"File saved successfully to {file_path}!")
                        try:
                            # Extract file extension to determine the type of processing required
                            _, file_extension = os.path.splitext(file_path)

                            text = extract_text_from_file(file_path, file_extension)
                
                            cleaned_file = clean_pdf_text(text)

                        except Exception as e:
                            print(f"Error creating directory: {e}")
                            logging.error(f"Error creating directory: {e}")
                    else:
                        print(f"Permission denied: Cannot write to {os.path.dirname(file_path)}")
                        logging.error(f"Permission denied: Cannot write to {os.path.dirname(file_path)}")
                        return f"Permission denied: Cannot write to {os.path.dirname(file_path)}"
                except Exception as save_error:
                    print(f"Error saving file: {save_error}")
                    logging.error(f"Error saving file: {save_error}")
                    error_message= "Oops, we had an issues saving your file. Please check your file and try again. Thanks for journeying with us!"
                    return StreamingResponse(error_stream(error_message), media_type='text/plain')
            else:
                print(f"Unsupported File: {file.filename}")
                logging.error(f"Unsupported File: {file.filename}")             
                error_message= "Oops, adventurous explorer! 🌍 You've ventured into uncharted territories. For now, let's stick to the well-trodden paths. Please try a supported format of pdf, txt, docx and please try not to use special characters. Thanks for journeying with us!"
                return StreamingResponse(error_stream(error_message), media_type='text/plain')
            
        try:
            if file:
                # If text bigger than the models max token then chunk and stream
                if num_tokens_from_string(cleaned_file, "cl100k_base") > 126000 and cleaned_file and cleaned_file.strip() != "":
                    print("Over 16000 tokens")
                    logging.info(f"FILE SIZE: FILE OVER File over 16000 tokens")
                    # Send the file content to ChatGPT

                    #summary_test_2 = "Summarize the uploaded annual report, focusing on. Detail Key business objectives and strategies for the upcoming year. Detail Major Initiatives and Projects. Detail Operational challenges and areas of concern. Outline who are the key decision-makers and what are their names and titles apart from the board. Is there any indication of digital transformation or areas where technology could provide a solution. Market Position and Competitive Landscape. Employee satisfaction, turnover rates, and growth. Customer Feedback and Trends. Regulatory and Compliance Concerns. Mergers and Acquisitions. Partnerships and Alliances"            
                    responses = send_text(prompt=user_request, text_data=cleaned_file)       

                    concatenated_responses = ' '.join(responses)
                    return StreamingResponse(custom_response_stream(concatenated_responses), media_type='text/plain')
                
                else:
                    print("Smaller")
                    logging.info(f"FILE SIZE: FILE UNDER File under 128000 tokens")

                    messages = [{"role": "system", "content": "You are a helpful assistant"}, {"role": "user", "content": f"Prompt: {user_request}\nContent:{cleaned_file}"}]

                    response_stream = openai.ChatCompletion.create(
                        model="gpt-4-1106-preview",
                        messages=messages,
                        stream=True
                    )#

                    # Use the event stream to send back the summary
                    def event_stream():
                        for line in response_stream:  
                            #print(line)
                            text = line.choices[0].delta.get('content', '')
                            #print(text)
                            yield text

                    return StreamingResponse(event_stream(), media_type='text/event-stream')
        except Exception as e:
            logging.error(f"An error occurred: {str(e)}")
            #error_message= "Oops, adventurous explorer! 🌍 You've ventured into uncharted territories. For now, let's stick to the well-trodden paths. Please try a supported format of pdf, txt, docx and please try not to use special characters. Thanks for journeying with us!"
            error_message= "Whoa there, speedster! 🚀 Looks like you're outpacing our preview release. Let's take a tiny breather. Please space out your questions a bit. Thanks for your support!"
            return StreamingResponse(error_stream(error_message), media_type='text/plain')
   

        #print("#######asdadasdadsas")
        get_news = await get_news_details(user_request)

        get_user_guide_data = await get_sf_user_guides_details(user_request)
          
        concatenated_string = get_news + get_user_guide_data

        final_query = (
            f"Based on the user's question: '[{user_request}]' and the accompanying supporting data reflecting key concepts and themes relevant to the users quesiton, create a structured summary of the following text. Focus on the primary points, supporting evidence, and conclusions. Ensure the summary is clear, accurate, and directly responsive to the user's query. Organize the content logically for a coherent and concise overview, enabling quick comprehension of the material."
            f"\nsupporting data: [{concatenated_string}]"
        )
        
        logging.info(f"FULL PROMPT: {final_query}")

        try:
            
            if num_tokens_from_string(final_query, "cl100k_base") > 126000:
                print("query longer than 16000")
                #summary_test_2 = "Summarize the uploaded annual report, focusing on. Detail Key business objectives and strategies for the upcoming year. Detail Major Initiatives and Projects. Detail Operational challenges and areas of concern. Outline who are the key decision-makers and what are their names and titles apart from the board. Is there any indication of digital transformation or areas where technology could provide a solution. Market Position and Competitive Landscape. Employee satisfaction, turnover rates, and growth. Customer Feedback and Trends. Regulatory and Compliance Concerns. Mergers and Acquisitions. Partnerships and Alliances"            
                responses = send_text(prompt=user_request, text_data=final_query)     
                #responses = send_text_local_model(prompt=summary_test_2, text_data=cleaned_file)       
                #send_text_local_model 
                # Use the event stream to send back the summary
                concatenated_responses = ' '.join(responses)
                return StreamingResponse(custom_response_stream(concatenated_responses), media_type='text/plain')
            
            else: 
                
                #messages_data = [{"role": "system", "content": ""}, {"role": "user", "content": f"{final_query}"}]
                messages_data = [{"role": "system", "content": "Pretend you are a expert assistant."}, {"role": "user", "content": f"{final_query}"}]

                messages = [Message(**message_data) for message_data in messages_data]

                #"gpt-3.5-turbo-16k or gpt-4"
                assistant_response = generate(messages, "gpt-4-1106-preview")   
                print("##################")

                return StreamingResponse(assistant_response, media_type='text/event-stream')
        except Exception as e:
            logging.error(f"An error occurred: {str(e)}")
            #error_message= "Oops, adventurous explorer! 🌍 You've ventured into uncharted territories. For now, let's stick to the well-trodden paths. Please try a supported format of pdf, txt, docx and please try not to use special characters. Thanks for journeying with us!"
            error_message= "Whoa there, speedster! 🚀 Looks like you're outpacing our preview release. Let's take a tiny breather. Please space out your questions a bit. Thanks for your support!"
            return StreamingResponse(error_stream(error_message), media_type='text/plain')
   
    except Exception as e:
        # Handling specific error codes
        logging.error(f"An error occurred: {str(e)}")
        error_message = "Whoa there, speedster! 🚀 The engine is currently overloaded with popular demand, please try again later."
        return StreamingResponse(error_stream(error_message), media_type='text/plain')




def get_current_user(session_token: str = Cookie(None)):
    if session_token is None or not r.exists(session_token):
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Not authenticated")
    return r.get(session_token).decode()  # Assuming the stored value is the email of the user


def create_session_token(email: str):
    session_token = ''.join(random.choice(string.ascii_letters + string.digits) for _ in range(12))
    # Store the session token in Redis with a 7-day expiration time
    r.setex(session_token, timedelta(days=7), value=email)
    return session_token


def store_credentials(email: str, password: str):
    hashed_password = pwd_context.hash(password)
    with open(CREDENTIALS_FILE, "a") as file:
        file.write(f"{email}:{hashed_password}\n")


def retrieve_hashed_password(email: str) -> str:
    if not os.path.exists(CREDENTIALS_FILE):
        return None
    with open(CREDENTIALS_FILE, "r") as file:
        lines = file.readlines()
        for line in lines:
            stored_email, hashed_password = line.strip().split(":")
            if stored_email == email:
                return hashed_password
    return None


def random_string(length=6):
    return ''.join(random.choice(string.ascii_uppercase + string.digits) for _ in range(length))

@app.get("/chat", response_class=HTMLResponse)
async def chat(request: Request, session_token: str = Cookie(None)):
    if session_token is None or not r.exists(session_token):
        # Redirect to the index page if the user is not authenticated
        return RedirectResponse(url="/", status_code=status.HTTP_303_SEE_OTHER)
    user_email = r.get(session_token).decode()  # Retrieve the user's email
    return templates.TemplateResponse("chat.html", {"request": request, "user": user_email})
 
@app.post("/signup")
async def signup(request: Request, email: str = Form(...)):
    print("Signup function is called.")  # Confirm the function is being called
    
    if not email.endswith("@salesforce.com"): 
        print("Email domain is invalid.")  # Log if the domain check fails
        return templates.TemplateResponse("index.html", {"request": request, "signup_error": "You require a Salesforce email to use this app."})
    
    # Check if the user already has an account
    if retrieve_hashed_password(email):
        print(f"Account already exists for {email}.")  # Log if the user already has an account
        return templates.TemplateResponse("index.html", {
            "request": request, 
            "signup_error": "An account with this email already exists. Please reset your password if you've forgotten it."
        })
    
    print(f"Attempting to create token for {email}")  # Confirm token creation
    token = random_string()
    r.set(f"signup_{email}", token, ex=3600)  # Setting the expiry time for 1 hour for signup verification

    print(f"Attempting to send verification email to {email}")  # Confirm email sending
    success = send_verification_email(email, token)

    if success:
        print(f"Verification email sent to {email}.")  # Log on success
        return templates.TemplateResponse("enter_token.html", {"request": request, "email": email})
    else:
        print(f"Failed to send verification email to {email}.")  # Log on failure
        return templates.TemplateResponse("index.html", {
            "request": request, 
            "signup_error": "Could not send verification email. Please try again later."
        })


@app.post("/verify-token")
async def verify_token(request: Request, email: str = Form(...), token: str = Form(...)):
    stored_token = r.get(f"signup_{email}")
    
    if not stored_token or stored_token.decode() != token:
        # Redirect back with an error if the token isn't found
        return templates.TemplateResponse("enter_token.html", {
            "request": request,
            "error": "Token not found or expired. Please try again.",  
            "email": email
        })
    else:
        # If everything is fine, proceed to set a new password
        return templates.TemplateResponse("set_password.html", {
            "request": request,
            "email": email,
            "token": token
        })


@app.post("/set-new-password")
async def set_new_password(request: Request, token: str = Form(...), email: str = Form(...), password: str = Form(...)):
    # Initialize response to None
    response = None

    # Verify the token, it could be a signup token or a reset token
    stored_signup_token = r.get(f"signup_{email}")
    stored_reset_token = r.get(f"reset_{email}")

    if stored_signup_token and stored_signup_token.decode() == token:
        # It's a sign-up process
        # Store the new password and remove the signup token
        store_credentials(email, password)
        r.delete(f"signup_{email}")
        # Create the session token and response
        session_token = create_session_token(email)
        response = RedirectResponse(url="/chat", status_code=status.HTTP_303_SEE_OTHER)
        response.set_cookie(key="session_token", value=session_token, max_age=604800)  # 7 days in seconds
    elif stored_reset_token and stored_reset_token.decode() == token:
        # It's a password reset process
        # Store the new password and remove the reset token
        store_credentials(email, password)
        r.delete(f"reset_{email}")
        # Create the session token and response
        session_token = create_session_token(email)
        response = RedirectResponse(url="/chat", status_code=status.HTTP_303_SEE_OTHER)
        response.set_cookie(key="session_token", value=session_token, max_age=604800)  # 7 days in seconds
    else:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid or expired token.")
    
    # If the response has not been created due to missing tokens, we should not proceed.
    if response is None:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Response could not be created.")
    
    return response



@app.post("/login") 
async def login(request: Request, email: str = Form(...), password: str = Form(...)):
    hashed_password = retrieve_hashed_password(email)
    if not hashed_password or not pwd_context.verify(password, hashed_password):
        # If the credentials are not valid, return the login page with an error
        return templates.TemplateResponse("index.html", {"request": request, "login_error": "Incorrect email or password"})
    
    # If the credentials are valid, create a session token and store it in Redis
    session_token = create_session_token(email)

    # Create the RedirectResponse with a 303 status code
    
    response = RedirectResponse(url="/chat", status_code=status.HTTP_303_SEE_OTHER)
    # Set the session_token in the cookie with max_age set to 7 days
    response.set_cookie(key="session_token", value=session_token, max_age=604800)  # 7 days in seconds
    return response


@app.get("/chat")
def main_page(request: Request, user_email: str = Depends(get_current_user)):

    # TODO: Implement session validation and render the main page
    # For now, let's just render a placeholder main page
    return templates.TemplateResponse("main_page.html", {"request": request})



@app.get("/reset-password-request")
async def reset_password_request(request: Request):
    return templates.TemplateResponse("reset_password_request.html", {"request": request})


@app.get("/reset-password")
async def reset_password_page(request: Request, token: str):
    email = None
    for key in r.scan_iter("reset_*"):
        if r.get(key).decode() == token:
            # Extract the email from the key
            email = key.decode()[6:]  # Assumes that key is of the form 'reset_{email}'
            break

    if email is None:
        # If the token is invalid or expired, inform the user
        return templates.TemplateResponse("reset_password.html", {
            "request": request,
            "error": "Invalid or expired token."
        })
    else:
        # If the token is valid, render the password reset form
        return templates.TemplateResponse("reset_password.html", {
            "request": request,
            "email": email,
            "token": token
        })


@app.post("/forgot-password")
async def forgot_password(request: Request, email: str = Form(...)):

        # Check if the email has the correct domain
    if not email.endswith("@salesforce.com"):
        # If the domain is incorrect, return the forgot password page with an error
        return templates.TemplateResponse("reset_password_request.html", {
            "request": request,
            "error": "Only Salesforce domain emails are allowed to reset passwords."
        })
    # Assuming email validation and user existence checks are done before this step.

    # Check if the user email exists in the system
    if not retrieve_hashed_password(email):
        print(f"No account exists for {email}.")  # Log if the user does not have an account
        return templates.TemplateResponse("reset_password_request.html", {
            "request": request, 
            "error": "No account found with this email. Please <a href='/'>signup</a> if you haven't already."
        })


    # Create a reset token and store it in Redis with a prefix 'reset_'
    reset_token = random_string()
    r.set(f"reset_{email}", reset_token, ex=600)  # Token expires after 600 seconds

    # Generate the reset link
    reset_link = f"http://127.0.0.1:8000/reset-password?token={reset_token}&email={email}"

    email_subject = "Password Reset Request"
    email_body = f"Please click on the link to reset your password: {reset_link}"
    
    # Use the send_email function
    success = send_email(email, email_subject, email_body)
    
    if success:
        # If the email was sent successfully, inform the user to check their email.
        return templates.TemplateResponse("check_your_email.html", {"request": request})
    else:
        # If there was an error sending the email, inform the user.
        return templates.TemplateResponse("reset_password_request.html", {
            "request": request,
            "error": "Could not send password reset email. Please try again."
        })


app.mount("/", StaticFiles(directory="static"), name="static")

if __name__ == '__main__':
    import uvicorn
    uvicorn.run(app, host='0.0.0.0', port=8000)

